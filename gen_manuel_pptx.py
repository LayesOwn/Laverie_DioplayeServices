"""Génère le manuel PowerPoint de Dioplaye Services."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x0A, 0x2D, 0x6E)
BLEU_VIF    = RGBColor(0x0D, 0x6E, 0xFD)
VERT        = RGBColor(0x19, 0x87, 0x54)
ROUGE       = RGBColor(0xDC, 0x35, 0x45)
ORANGE      = RGBColor(0xFF, 0xC1, 0x07)
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR  = RGBColor(0xF0, 0xF4, 0xF8)
GRIS_TEXTE  = RGBColor(0x49, 0x54, 0x5E)
NOIR        = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # layout vide


def add_rect(slide, l, t, w, h, fill=None, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=NOIR, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, l, t, w, h, size=15, color=GRIS_TEXTE,
               bullet="●", title=None, title_color=BLEU_VIF):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, l, t, w, h, bg, fg=BLANC, size=13):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, text, l, t, w, h, size=size, bold=True,
             color=fg, align=PP_ALIGN.CENTER)


def slide_titre(prs):
    """Slide 1 — Titre"""
    sl = prs.slides.add_slide(BLANK)

    # Fond dégradé simulé : deux rectangles
    add_rect(sl, 0, 0, W, H, fill=BLEU_FONCE)
    add_rect(sl, 0, H * 0.55, W, H * 0.45, fill=BLEU_VIF)

    # Bande décorative
    add_rect(sl, 0, Inches(4.2), W, Inches(0.06), fill=BLANC)

    # Logo texte
    add_text(sl, "DIOPLAYE SERVICES", Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             size=54, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    add_text(sl, "Manuel d'utilisation", Inches(1), Inches(2.6), Inches(11), Inches(0.8),
             size=28, bold=False, color=RGBColor(0xBD, 0xD7, 0xFF), align=PP_ALIGN.CENTER,
             italic=True)

    add_text(sl, "Application de gestion de laverie", Inches(1), Inches(3.3),
             Inches(11), Inches(0.6),
             size=18, color=RGBColor(0x90, 0xB8, 0xFF), align=PP_ALIGN.CENTER)

    # Badges
    for i, (txt, bg) in enumerate([
        ("Gestion Clients", VERT),
        ("Transactions", BLEU_VIF),
        ("Exports CSV / Excel / PDF", ORANGE),
        ("Application Mobile (PWA)", ROUGE),
    ]):
        add_badge(sl, txt, Inches(1 + i * 2.85), Inches(5.2),
                  Inches(2.6), Inches(0.5), bg=bg, size=12)

    add_text(sl, "https://layes03.pythonanywhere.com",
             Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             size=13, color=RGBColor(0xBD, 0xD7, 0xFF), align=PP_ALIGN.CENTER)


def header(sl, titre, subtitle=None):
    add_rect(sl, 0, 0, W, Inches(1.1), fill=BLEU_FONCE)
    add_rect(sl, 0, Inches(1.1), W, Inches(0.05), fill=BLEU_VIF)
    add_text(sl, titre, Inches(0.4), Inches(0.08), Inches(9), Inches(0.7),
             size=26, bold=True, color=BLANC)
    if subtitle:
        add_text(sl, subtitle, Inches(0.4), Inches(0.72), Inches(12), Inches(0.38),
                 size=13, color=RGBColor(0xBD, 0xD7, 0xFF))
    # numéro de slide (approximatif)
    add_text(sl, "Dioplaye Services", Inches(10.5), Inches(0.3),
             Inches(2.5), Inches(0.5), size=11,
             color=RGBColor(0x90, 0xB8, 0xFF), align=PP_ALIGN.RIGHT)


def card(sl, titre, corps, l, t, w, h, accent=BLEU_VIF):
    add_rect(sl, l, t, w, h, fill=GRIS_CLAIR, line=accent)
    add_rect(sl, l, t, Inches(0.07), h, fill=accent)
    add_text(sl, titre, l + Inches(0.15), t + Inches(0.1),
             w - Inches(0.2), Inches(0.38),
             size=13, bold=True, color=accent)
    add_text(sl, corps, l + Inches(0.15), t + Inches(0.5),
             w - Inches(0.2), h - Inches(0.6),
             size=12, color=GRIS_TEXTE)


# ─────────────────────────────────────────────────────────────
# SLIDE 1 — Titre
slide_titre(prs)


# ─────────────────────────────────────────────────────────────
# SLIDE 2 — Présentation générale
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Présentation de l'application",
       "Dioplaye Services — Gestion complète d'une laverie")

add_text(sl, "Qu'est-ce que c'est ?", Inches(0.5), Inches(1.3),
         Inches(12), Inches(0.5), size=18, bold=True, color=BLEU_FONCE)

add_text(sl,
         "Dioplaye Services est une application web de gestion complète pour laverie. "
         "Elle permet de suivre les clients, les transactions, les paiements, "
         "les dépenses et de générer des rapports et reçus.",
         Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.9),
         size=14, color=GRIS_TEXTE)

fonctionnalites = [
    ("Gestion des clients", "Fiches clients, historique, solde dû en temps réel", BLEU_VIF),
    ("Transactions & Paiements", "Enregistrement des commandes et encaissements", VERT),
    ("Dépenses internes", "Suivi des charges par catégorie (lessive, électricité…)", ORANGE),
    ("Tableau de bord", "KPIs du jour, du mois, de l'année + graphiques", BLEU_FONCE),
    ("Exports", "CSV, Excel et reçus PDF imprimables", ROUGE),
    ("Application mobile", "Installable sur téléphone Android via Chrome (PWA)", VERT),
]
cols = [(0.4, 4.0), (4.65, 4.0), (9.0, 4.0)]
for i, (titre_c, corps_c, acc) in enumerate(fonctionnalites):
    col_i = i % 3
    row_i = i // 3
    l = Inches(cols[col_i][0])
    w = Inches(cols[col_i][1])
    t = Inches(2.9 + row_i * 1.7)
    card(sl, titre_c, corps_c, l, t, w, Inches(1.55), accent=acc)


# ─────────────────────────────────────────────────────────────
# SLIDE 3 — Connexion et rôles
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Connexion & Rôles utilisateurs",
       "Deux niveaux d'accès sécurisés")

# Schéma connexion
add_rect(sl, Inches(0.4), Inches(1.3), Inches(12.4), Inches(1.0), fill=GRIS_CLAIR, line=BLEU_VIF)
add_text(sl, "URL : https://layes03.pythonanywhere.com  →  Page de connexion (nom d'utilisateur + mot de passe)",
         Inches(0.6), Inches(1.48), Inches(12), Inches(0.6),
         size=14, color=BLEU_FONCE, bold=False)

# Carte Admin
add_rect(sl, Inches(0.4), Inches(2.6), Inches(5.8), Inches(4.3), fill=RGBColor(0xE8, 0xF4, 0xFF), line=BLEU_VIF)
add_rect(sl, Inches(0.4), Inches(2.6), Inches(5.8), Inches(0.55), fill=BLEU_VIF)
add_text(sl, "ADMINISTRATEUR", Inches(0.5), Inches(2.65), Inches(5.6), Inches(0.45),
         size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

admin_droits = [
    "Toutes les fonctions de l'invité",
    "Supprimer clients, transactions, dépenses",
    "Gérer les services (créer, modifier, prix)",
    "Accéder à la corbeille et restaurer",
    "Créer des comptes invités",
    "Consulter les recettes historiques",
    "Accès complet aux exports",
]
bullet_box(sl, admin_droits, Inches(0.6), Inches(3.25), Inches(5.4), Inches(3.5),
           size=13, color=GRIS_TEXTE)

# Carte Invité
add_rect(sl, Inches(6.6), Inches(2.6), Inches(6.2), Inches(4.3), fill=RGBColor(0xE8, 0xFF, 0xF0), line=VERT)
add_rect(sl, Inches(6.6), Inches(2.6), Inches(6.2), Inches(0.55), fill=VERT)
add_text(sl, "INVITÉ (Employé)", Inches(6.7), Inches(2.65), Inches(6.0), Inches(0.45),
         size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

invite_droits = [
    "Créer et modifier des clients",
    "Enregistrer des transactions",
    "Encaisser des paiements",
    "Ajouter des dépenses",
    "Consulter tous les clients et transactions",
    "Télécharger les exports",
]
bullet_box(sl, invite_droits, Inches(6.8), Inches(3.25), Inches(5.8), Inches(3.5),
           size=13, color=GRIS_TEXTE)


# ─────────────────────────────────────────────────────────────
# SLIDE 4 — Tableau de bord
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Tableau de bord",
       "Vue d'ensemble en temps réel de l'activité")

kpis = [
    ("Recettes du jour",    "Paiements reçus aujourd'hui",        VERT),
    ("Dépenses du jour",    "Charges enregistrées aujourd'hui",   ROUGE),
    ("Bénéfice net / jour", "Recettes − Dépenses du jour",        BLEU_VIF),
    ("Recettes du mois",    "Cumul du mois en cours",             VERT),
    ("Dépenses du mois",    "Charges du mois en cours",           ROUGE),
    ("Bénéfice mensuel",    "Recettes − Dépenses du mois",        BLEU_VIF),
]
for i, (titre_k, desc_k, acc_k) in enumerate(kpis):
    col = i % 3
    row = i // 3
    l = Inches(0.35 + col * 4.3)
    t = Inches(1.3 + row * 1.35)
    add_rect(sl, l, t, Inches(4.0), Inches(1.2), fill=GRIS_CLAIR, line=acc_k)
    add_rect(sl, l, t, Inches(0.08), Inches(1.2), fill=acc_k)
    add_text(sl, titre_k, l + Inches(0.15), t + Inches(0.1),
             Inches(3.7), Inches(0.45), size=13, bold=True, color=acc_k)
    add_text(sl, desc_k, l + Inches(0.15), t + Inches(0.55),
             Inches(3.7), Inches(0.55), size=12, color=GRIS_TEXTE)

add_rect(sl, Inches(0.35), Inches(4.1), Inches(12.6), Inches(0.05), fill=BLEU_VIF)

graphiques = [
    "Évolution mensuelle sur 12 mois (recettes, dépenses, bénéfice)",
    "Répartition des services les plus demandés (camembert)",
    "Recettes par mois par année (courbes historiques)",
    "Top 5 clients (plus gros payeurs)",
    "Top 5 services (les plus utilisés)",
    "Créances : nombre et montant total dû par les clients",
]
bullet_box(sl, graphiques, Inches(0.5), Inches(4.25), Inches(12), Inches(3.0),
           size=13, color=GRIS_TEXTE,
           title="Graphiques et indicateurs inclus :", title_color=BLEU_FONCE)


# ─────────────────────────────────────────────────────────────
# SLIDE 5 — Clients
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Gestion des clients",
       "Fiches clients, historique complet, suivi des impayés")

col1 = [
    ("Ajouter un client", [
        "Nom complet (obligatoire)",
        "Téléphone (recommandé)",
        "Adresse (optionnel)",
        "Remarque personnalisée",
        "Première transaction optionnelle au même formulaire",
    ], BLEU_VIF),
    ("Rechercher un client", [
        "Barre de recherche par nom",
        "Résultats en temps réel",
        "Pagination (15 par page)",
    ], BLEU_FONCE),
]
col2 = [
    ("Fiche client détaillée", [
        "Total facturé sur toutes les transactions",
        "Total déjà payé",
        "Solde restant dû (en rouge si impayé)",
        "Historique complet de toutes les transactions",
    ], VERT),
    ("Modifier / Supprimer", [
        "Modifier : nom, téléphone, adresse, remarque",
        "Supprimer (admin) : efface le client et tout son historique",
    ], ROUGE),
]

for row_i, (titre_c, items_c, acc_c) in enumerate(col1):
    t = Inches(1.35 + row_i * 2.8)
    add_rect(sl, Inches(0.3), t, Inches(6.2), Inches(2.6), fill=GRIS_CLAIR, line=acc_c)
    add_rect(sl, Inches(0.3), t, Inches(0.1), Inches(2.6), fill=acc_c)
    add_text(sl, titre_c, Inches(0.55), t + Inches(0.1), Inches(5.8), Inches(0.45),
             size=14, bold=True, color=acc_c)
    bullet_box(sl, items_c, Inches(0.55), t + Inches(0.55), Inches(5.8), Inches(1.9),
               size=12, color=GRIS_TEXTE)

for row_i, (titre_c, items_c, acc_c) in enumerate(col2):
    t = Inches(1.35 + row_i * 2.8)
    add_rect(sl, Inches(6.8), t, Inches(6.2), Inches(2.6), fill=GRIS_CLAIR, line=acc_c)
    add_rect(sl, Inches(6.8), t, Inches(0.1), Inches(2.6), fill=acc_c)
    add_text(sl, titre_c, Inches(7.05), t + Inches(0.1), Inches(5.8), Inches(0.45),
             size=14, bold=True, color=acc_c)
    bullet_box(sl, items_c, Inches(7.05), t + Inches(0.55), Inches(5.8), Inches(1.9),
               size=12, color=GRIS_TEXTE)


# ─────────────────────────────────────────────────────────────
# SLIDE 6 — Transactions
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Transactions",
       "Enregistrement des commandes clients avec gestion du crédit")

etapes = [
    ("1. Choisir le client",
     "Sélectionner dans la liste  OU  rechercher par nom/téléphone  OU  créer un nouveau client.\n"
     "L'app détecte automatiquement un client existant (même nom ou téléphone)."),
    ("2. Choisir le service",
     "Sélectionner la prestation. Le prix unitaire s'affiche automatiquement.\n"
     "Saisir la quantité (ex: 2 pour 2 machines). Le total se calcule seul."),
    ("3. Saisir l'avance",
     "Montant payé immédiatement à la commande (peut être 0).\n"
     "Le reste dû est calculé et visible sur la fiche client."),
    ("4. Date et notes",
     "Date de la transaction (aujourd'hui par défaut).\n"
     "Notes libres (ex: \"chemises blanches, programme délicat\")."),
]

for i, (titre_e, corps_e) in enumerate(etapes):
    t = Inches(1.3 + i * 1.45)
    add_rect(sl, Inches(0.3), t, Inches(0.7), Inches(1.25), fill=BLEU_VIF)
    add_text(sl, str(i + 1), Inches(0.3), t + Inches(0.25),
             Inches(0.7), Inches(0.7), size=28, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.1), t, Inches(11.8), Inches(1.25), fill=GRIS_CLAIR, line=BLEU_VIF)
    add_text(sl, titre_e, Inches(1.25), t + Inches(0.08),
             Inches(11.4), Inches(0.4), size=14, bold=True, color=BLEU_FONCE)
    add_text(sl, corps_e, Inches(1.25), t + Inches(0.5),
             Inches(11.4), Inches(0.7), size=12, color=GRIS_TEXTE)

add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32), fill=RGBColor(0xFF, 0xF3, 0xCD))
add_rect(sl, Inches(0.3), Inches(7.1), Inches(0.07), Inches(0.32), fill=ORANGE)
add_text(sl, "Filtre \"Non soldées\" dans la liste des transactions pour voir uniquement les clients qui doivent encore de l'argent.",
         Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.28),
         size=11, color=RGBColor(0x85, 0x60, 0x04))


# ─────────────────────────────────────────────────────────────
# SLIDE 7 — Paiements
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Enregistrement des paiements",
       "Encaissement du solde restant au retrait du linge")

add_text(sl, "Comment encaisser un paiement ?",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         size=18, bold=True, color=BLEU_FONCE)

steps = [
    ("Ouvrir la fiche client", "Menu Clients → cliquer sur le nom du client"),
    ("Trouver la transaction", "Voir la liste des transactions avec le solde restant dû"),
    ("Cliquer Enregistrer paiement", "Bouton vert sur la ligne de la transaction concernée"),
    ("Saisir le montant reçu", "En XOF, entier (pas de décimales). L'app plafonne au solde dû."),
    ("Valider", "La transaction passe au statut SOLDÉE si le solde atteint 0"),
]

for i, (titre_s, corps_s) in enumerate(steps):
    t = Inches(1.95 + i * 0.95)
    add_rect(sl, Inches(0.3), t, Inches(0.5), Inches(0.78), fill=VERT)
    add_text(sl, "✓", Inches(0.3), t + Inches(0.1),
             Inches(0.5), Inches(0.55), size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.9), t, Inches(11.8), Inches(0.78), fill=GRIS_CLAIR, line=VERT)
    add_text(sl, titre_s, Inches(1.05), t + Inches(0.05),
             Inches(4.5), Inches(0.38), size=13, bold=True, color=VERT)
    add_text(sl, corps_s, Inches(1.05), t + Inches(0.42),
             Inches(11.4), Inches(0.32), size=12, color=GRIS_TEXTE)

add_rect(sl, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.45), fill=RGBColor(0xD1, 0xFA, 0xE5))
add_rect(sl, Inches(0.3), Inches(6.9), Inches(0.07), Inches(0.45), fill=VERT)
add_text(sl,
         "Protection anti-doublon : l'app vérifie en temps réel que le total ne dépasse "
         "jamais le montant facturé. Impossible de trop encaisser.",
         Inches(0.5), Inches(6.93), Inches(12.3), Inches(0.38),
         size=12, color=RGBColor(0x06, 0x5F, 0x46))


# ─────────────────────────────────────────────────────────────
# SLIDE 8 — Dépenses internes
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Dépenses internes",
       "Suivi des charges de la laverie par catégorie")

categories = [
    ("Lessive / Produits", BLEU_VIF),
    ("Électricité / Eau", ORANGE),
    ("Matériel / Équipement", BLEU_FONCE),
    ("Loyer", ROUGE),
    ("Transport", VERT),
    ("Salaires", RGBColor(0x6F, 0x42, 0xC1)),
    ("Maintenance", RGBColor(0x0D, 0xCA, 0xF0)),
    ("Autre", GRIS_TEXTE),
]

add_text(sl, "Catégories disponibles :", Inches(0.5), Inches(1.3),
         Inches(5), Inches(0.4), size=14, bold=True, color=BLEU_FONCE)

for i, (cat, acc) in enumerate(categories):
    col = i % 2
    row = i // 2
    l = Inches(0.4 + col * 2.6)
    t = Inches(1.8 + row * 0.6)
    add_rect(sl, l, t, Inches(2.3), Inches(0.45), fill=GRIS_CLAIR, line=acc)
    add_rect(sl, l, t, Inches(0.06), Inches(0.45), fill=acc)
    add_text(sl, cat, l + Inches(0.12), t + Inches(0.05),
             Inches(2.1), Inches(0.35), size=12, color=GRIS_TEXTE)

add_rect(sl, Inches(5.6), Inches(1.3), Inches(7.5), Inches(5.8), fill=GRIS_CLAIR, line=BLEU_VIF)
add_rect(sl, Inches(5.6), Inches(1.3), Inches(7.5), Inches(0.5), fill=BLEU_VIF)
add_text(sl, "Formulaire de saisie", Inches(5.75), Inches(1.35),
         Inches(7.2), Inches(0.4), size=14, bold=True, color=BLANC)

champs = [
    ("Libellé", "Description de la dépense  (ex: Ariel 5kg)"),
    ("Montant", "En XOF (chiffre entier)"),
    ("Type", "Dépense  ou  Achat"),
    ("Catégorie", "Choisir dans la liste ci-contre"),
    ("Date", "Date de la dépense"),
    ("Notes", "Informations complémentaires (optionnel)"),
]
for i, (champ, desc) in enumerate(champs):
    t = Inches(2.0 + i * 0.75)
    add_text(sl, champ, Inches(5.8), t, Inches(1.8), Inches(0.4),
             size=12, bold=True, color=BLEU_FONCE)
    add_text(sl, desc, Inches(7.65), t, Inches(5.2), Inches(0.4),
             size=12, color=GRIS_TEXTE)
    if i < len(champs) - 1:
        add_rect(sl, Inches(5.75), t + Inches(0.42), Inches(7.2), Inches(0.02), fill=RGBColor(0xDE, 0xE2, 0xE6))

add_rect(sl, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.42), fill=RGBColor(0xFF, 0xF3, 0xCD))
add_rect(sl, Inches(0.4), Inches(6.55), Inches(0.07), Inches(0.42), fill=ORANGE)
add_text(sl,
         "Suppression douce : les dépenses supprimées vont dans la corbeille (admin). "
         "Elles peuvent être restaurées à tout moment.",
         Inches(0.6), Inches(6.58), Inches(12.1), Inches(0.35),
         size=11, color=RGBColor(0x85, 0x60, 0x04))


# ─────────────────────────────────────────────────────────────
# SLIDE 9 — Services
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Gestion des services",
       "Prestations proposées aux clients — accès administrateur")

services_defaut = [
    ("Lavage simple", "1 500 XOF", "Lavage en machine sans séchage"),
    ("Lavage + Séchage", "2 500 XOF", "Lavage et séchage complet"),
    ("Lavage + Séchage + Repassage", "4 000 XOF", "Service complet"),
]

add_text(sl, "Services par défaut :", Inches(0.5), Inches(1.3),
         Inches(12), Inches(0.4), size=16, bold=True, color=BLEU_FONCE)

for i, (nom_s, prix_s, desc_s) in enumerate(services_defaut):
    t = Inches(1.85 + i * 1.2)
    add_rect(sl, Inches(0.4), t, Inches(8.5), Inches(1.05), fill=GRIS_CLAIR, line=BLEU_VIF)
    add_rect(sl, Inches(0.4), t, Inches(0.1), Inches(1.05), fill=BLEU_VIF)
    add_text(sl, nom_s, Inches(0.65), t + Inches(0.08), Inches(5), Inches(0.42),
             size=14, bold=True, color=BLEU_FONCE)
    add_text(sl, desc_s, Inches(0.65), t + Inches(0.54), Inches(5), Inches(0.42),
             size=12, color=GRIS_TEXTE)
    add_rect(sl, Inches(7.5), t, Inches(1.4), Inches(1.05), fill=BLEU_VIF)
    add_text(sl, prix_s, Inches(7.5), t + Inches(0.28),
             Inches(1.4), Inches(0.5), size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(9.3), Inches(1.3), Inches(3.8), Inches(5.8), fill=GRIS_CLAIR, line=BLEU_FONCE)
add_rect(sl, Inches(9.3), Inches(1.3), Inches(3.8), Inches(0.5), fill=BLEU_FONCE)
add_text(sl, "Actions disponibles (admin)", Inches(9.45), Inches(1.35),
         Inches(3.6), Inches(0.4), size=13, bold=True, color=BLANC)

actions = [
    "Créer un nouveau service",
    "Modifier le nom du service",
    "Modifier le prix unitaire",
    "Modifier la description",
    "Désactiver un service (n'apparaît plus dans les nouveaux formulaires)",
    "Supprimer un service inutilisé",
]
bullet_box(sl, actions, Inches(9.45), Inches(1.95), Inches(3.6), Inches(4.8),
           size=12, color=GRIS_TEXTE)

add_rect(sl, Inches(0.4), Inches(6.55), Inches(8.7), Inches(0.42), fill=RGBColor(0xD1, 0xFA, 0xE5))
add_rect(sl, Inches(0.4), Inches(6.55), Inches(0.07), Inches(0.42), fill=VERT)
add_text(sl,
         "Modifier un prix n'affecte pas les transactions déjà enregistrées.",
         Inches(0.6), Inches(6.58), Inches(8.3), Inches(0.35),
         size=12, color=RGBColor(0x06, 0x5F, 0x46))


# ─────────────────────────────────────────────────────────────
# SLIDE 10 — Exports
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Exports et rapports",
       "Téléchargement des données et impression des reçus")

exports_data = [
    ("Export CSV\nTransactions", BLEU_VIF,
     ["Toutes les transactions actives",
      "Colonnes : ID, Date, Client, Service,",
      "Quantité, Total, Payé, Reste dû, Notes",
      "Format : fichier .csv (Excel, Sheets)",
      "Menu sidebar → Export CSV"]),
    ("Export Excel\nTransactions", VERT,
     ["Même données que le CSV",
      "Format .xlsx avec mise en forme",
      "En-têtes colorés en bleu",
      "Largeurs de colonnes automatiques",
      "Menu sidebar → Export Excel"]),
    ("Export CSV\nDépenses", ORANGE,
     ["Toutes les dépenses actives",
      "Colonnes : Date, Libellé, Type,",
      "Catégorie, Montant, Notes",
      "Exclut les dépenses en corbeille",
      "Menu Exports → Export CSV dépenses"]),
    ("Reçu PDF\nClient", ROUGE,
     ["Un reçu par transaction",
      "Format A5 (demi-page A4)",
      "Contient : client, service, total,",
      "montant payé, reste dû, date",
      "Fiche transaction → bouton Reçu PDF"]),
]

for i, (titre_e, acc_e, items_e) in enumerate(exports_data):
    l = Inches(0.3 + i * 3.2)
    add_rect(sl, l, Inches(1.3), Inches(3.0), Inches(5.9), fill=GRIS_CLAIR, line=acc_e)
    add_rect(sl, l, Inches(1.3), Inches(3.0), Inches(0.65), fill=acc_e)
    add_text(sl, titre_e, l + Inches(0.1), Inches(1.33),
             Inches(2.8), Inches(0.6), size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    bullet_box(sl, items_e, l + Inches(0.15), Inches(2.05),
               Inches(2.8), Inches(4.9), size=12, color=GRIS_TEXTE, bullet="→")


# ─────────────────────────────────────────────────────────────
# SLIDE 11 — Corbeille
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Corbeille — Protection des données",
       "Les suppressions sont réversibles (admin uniquement)")

add_text(sl,
         "Aucune donnée n'est définitivement perdue lors d'une suppression. "
         "Les éléments supprimés sont déplacés dans une corbeille et peuvent être restaurés à tout moment.",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7),
         size=14, color=GRIS_TEXTE)

for i, (titre_c, items_c, acc_c, chemin) in enumerate([
    ("Corbeille Transactions", [
        "Transactions supprimées par erreur",
        "Chaque ligne affiche la date de suppression",
        "Bouton Restaurer pour réactiver",
        "Les paiements associés sont conservés",
    ], BLEU_VIF, "Menu Transactions → lien Corbeille"),
    ("Corbeille Dépenses", [
        "Dépenses supprimées par erreur",
        "Chaque ligne affiche la date de suppression",
        "Bouton Restaurer pour réactiver",
        "Dépense restaurée réapparaît dans la liste normale",
    ], ROUGE, "Menu Dépenses → lien Corbeille"),
]):
    l = Inches(0.4 + i * 6.6)
    add_rect(sl, l, Inches(2.15), Inches(6.2), Inches(4.2), fill=GRIS_CLAIR, line=acc_c)
    add_rect(sl, l, Inches(2.15), Inches(6.2), Inches(0.55), fill=acc_c)
    add_text(sl, titre_c, l + Inches(0.15), Inches(2.19),
             Inches(5.9), Inches(0.45), size=15, bold=True, color=BLANC)
    bullet_box(sl, items_c, l + Inches(0.2), Inches(2.8), Inches(5.8), Inches(2.8),
               size=13, color=GRIS_TEXTE)
    add_rect(sl, l, Inches(5.1), Inches(6.2), Inches(0.08), fill=acc_c)
    add_text(sl, f"Accès : {chemin}", l + Inches(0.15), Inches(5.2),
             Inches(5.9), Inches(0.4), size=12, bold=True, color=acc_c)

add_rect(sl, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.42), fill=RGBColor(0xFE, 0xE2, 0xE2))
add_rect(sl, Inches(0.4), Inches(6.55), Inches(0.07), Inches(0.42), fill=ROUGE)
add_text(sl,
         "Important : la suppression définitive depuis la corbeille n'est pas disponible "
         "(protection des données). Pour toute suppression permanente, contacter l'administrateur technique.",
         Inches(0.6), Inches(6.58), Inches(12.1), Inches(0.35),
         size=11, color=RGBColor(0x99, 0x1B, 0x1B))


# ─────────────────────────────────────────────────────────────
# SLIDE 12 — Installation PWA
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLANC)
header(sl, "Installer l'app sur téléphone (PWA)",
       "L'application fonctionne comme une vraie appli Android")

add_text(sl, "Qu'est-ce que la PWA ?", Inches(0.5), Inches(1.3),
         Inches(12), Inches(0.45), size=16, bold=True, color=BLEU_FONCE)
add_text(sl,
         "PWA = Progressive Web App. L'application s'installe depuis Chrome sur Android "
         "et s'ouvre sans barre de navigation, comme une application native.",
         Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.55),
         size=13, color=GRIS_TEXTE)

etapes_pwa = [
    ("Ouvrir Chrome", "sur ton téléphone Android"),
    ("Aller sur", "https://layes03.pythonanywhere.com"),
    ("Se connecter", "avec ton nom d'utilisateur et mot de passe"),
    ("Appuyer sur ⋮", "les 3 points en haut à droite de Chrome"),
    ("Choisir", '"Ajouter à l\'écran d\'accueil"'),
    ("Confirmer", "L'icône bleue apparaît sur l'écran d'accueil"),
]

for i, (action, detail) in enumerate(etapes_pwa):
    col = i % 3
    row = i // 2
    l = Inches(0.3 + col * 4.35)
    t = Inches(2.6 + row * 1.55)
    add_rect(sl, l, t, Inches(4.0), Inches(1.35), fill=GRIS_CLAIR, line=BLEU_VIF)
    add_rect(sl, l, t, Inches(0.55), Inches(1.35), fill=BLEU_VIF)
    add_text(sl, str(i + 1), l + Inches(0.0), t + Inches(0.35),
             Inches(0.55), Inches(0.65), size=22, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(sl, action, l + Inches(0.65), t + Inches(0.15),
             Inches(3.2), Inches(0.4), size=13, bold=True, color=BLEU_FONCE)
    add_text(sl, detail, l + Inches(0.65), t + Inches(0.6),
             Inches(3.2), Inches(0.65), size=12, color=GRIS_TEXTE)

avantages = [
    "Icône sur l'écran d'accueil",
    "Ouverture plein écran (sans barre Chrome)",
    "Session maintenue (pas de reconnexion à chaque fois)",
    "Fonctionne sur réseau mobile et WiFi",
]
bullet_box(sl, avantages, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.95),
           size=12, color=GRIS_TEXTE,
           title="Avantages :  ", title_color=VERT)


# ─────────────────────────────────────────────────────────────
# SLIDE 13 — Conclusion
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=BLEU_FONCE)
add_rect(sl, 0, Inches(5.8), W, Inches(1.7), fill=BLEU_VIF)
add_rect(sl, 0, Inches(4.6), W, Inches(0.06), fill=BLANC)

add_text(sl, "Dioplaye Services", Inches(1), Inches(0.8), Inches(11), Inches(1.0),
         size=42, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text(sl, "Votre laverie, gérée simplement.", Inches(1), Inches(1.95), Inches(11), Inches(0.6),
         size=22, color=RGBColor(0xBD, 0xD7, 0xFF), align=PP_ALIGN.CENTER, italic=True)

recap = [
    ("Clients & Transactions",  BLEU_VIF),
    ("Paiements & Créances",    VERT),
    ("Dépenses par catégorie",  ORANGE),
    ("Exports CSV / Excel / PDF", ROUGE),
    ("Tableau de bord & Graphiques", RGBColor(0x0D, 0xCA, 0xF0)),
    ("Application mobile (PWA)", RGBColor(0x6F, 0x42, 0xC1)),
]
for i, (txt_r, acc_r) in enumerate(recap):
    col = i % 3
    row = i // 3
    l = Inches(0.5 + col * 4.3)
    t = Inches(2.8 + row * 0.85)
    add_rect(sl, l, t, Inches(4.0), Inches(0.65), fill=acc_r)
    add_text(sl, "✓  " + txt_r, l + Inches(0.1), t + Inches(0.1),
             Inches(3.8), Inches(0.45), size=13, bold=True, color=BLANC)

add_text(sl, "https://layes03.pythonanywhere.com",
         Inches(1), Inches(5.95), Inches(11), Inches(0.45),
         size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text(sl, "Support : dioplayes@gmail.com",
         Inches(1), Inches(6.5), Inches(11), Inches(0.4),
         size=13, color=RGBColor(0xBD, 0xD7, 0xFF), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# Sauvegarde
output = "Manuel_DioplayeServices.pptx"
prs.save(output)
print(f"Fichier cree : {output}")
